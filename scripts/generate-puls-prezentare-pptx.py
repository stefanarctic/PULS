# -*- coding: utf-8 -*-
"""
Prezentare P.U.L.S. — layout modern (bandă titlu sombre, corp deschis, cadru foto),
capturi din `public/res/screenshots`, tranziție Morph injectată OOXML după salvare.

Fișier: docs/PULS-Prezentare-site.pptx
"""
from __future__ import annotations

import io
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "PULS-Prezentare-site.pptx"
SHOTS = ROOT / "public" / "res" / "screenshots"
LOGO = ROOT / "public" / "res" / "icons" / "New-logo.png"

SW = Inches(13.333)
SH = Inches(7.5)

C_BG = RGBColor(248, 250, 252)
C_BAR = RGBColor(30, 41, 59)
C_TEXT = RGBColor(51, 65, 85)
C_LINK = RGBColor(79, 70, 229)
C_FRAME = RGBColor(148, 163, 184)


def resolve_img(*candidates: str) -> Path:
    for name in candidates:
        p = SHOTS / name
        if p.is_file():
            return p
    if LOGO.is_file():
        return LOGO
    raise FileNotFoundError(f"Lipsă imagini {candidates} și logo {LOGO}")


def _no_line(shape) -> None:
    shape.line.fill.background()


def inject_morph(ppx_src: Path, ppx_dest: Path) -> None:
    morph_snippet = (
        '<p:transition spd="med">'
        '<p14:morph option="true" xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"/>'
        "</p:transition>"
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(ppx_src, "r") as zin, zipfile.ZipFile(
        buf, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            fn = item.filename
            if (
                fn.startswith("ppt/slides/slide")
                and fn.endswith(".xml")
                and "/_rels/" not in fn
            ):
                text = data.decode("utf-8")
                if "<p:transition" not in text and "</p:sld>" in text:
                    text = text.replace("</p:sld>", morph_snippet + "</p:sld>", 1)
                data = text.encode("utf-8")
            zout.writestr(item, data)
    buf.seek(0)
    ppx_dest.write_bytes(buf.read())


def save_presentation_with_morph(prs: Presentation, target: Path) -> None:
    staging = target.with_suffix(".staging.pptx")
    try:
        prs.save(staging)
        inject_morph(staging, target)
    finally:
        if staging.exists():
            staging.unlink(missing_ok=True)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _no_line(bg)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.name = "BGFill"

    bar_h = Inches(1.62)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, bar_h)
    _no_line(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_BAR
    bar.name = "TitleBandTop"

    if LOGO.is_file():
        lg = slide.shapes.add_picture(
            str(LOGO), Inches(0.5), Inches(0.32), height=Inches(0.92)
        )
        lg.name = "BrandLogo"

    title_box = slide.shapes.add_textbox(
        Inches(1.18), Inches(0.22), SW - Inches(1.4), Inches(1.15)
    )
    title_box.name = "DeckTitleMain"
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.clear()
    ta = tf.paragraphs[0]
    ta.text = "P.U.L.S."
    ta.font.size = Pt(54)
    ta.font.bold = True
    ta.font.color.rgb = RGBColor(255, 255, 255)

    tb = tf.add_paragraph()
    tb.text = "Simulări · Bac · Resurse · Asistent AI — în același loc"
    tb.font.size = Pt(13)
    tb.font.color.rgb = RGBColor(203, 213, 245)
    tb.space_before = Pt(2)

    tc = tf.add_paragraph()
    tc.text = "https://puls-fizica.ro"
    tc.font.size = Pt(14)
    tc.font.bold = True
    tc.font.color.rgb = RGBColor(196, 181, 253)
    tc.space_before = Pt(12)

    hero = resolve_img(
        "Laser_Simulator_Screenshot.png",
        "Termodinamica_Screenshot.png",
        "Unde_Screenshot.png",
        "Prisma_Screenshot.png",
    )
    pic_top = bar_h + Inches(0.22)
    ph = slide.shapes.add_picture(
        str(hero), Inches(0.42), pic_top, width=SW - Inches(0.84)
    )
    ph.name = "TitleHeroShot"
    max_photo_h = SH - pic_top - Inches(0.38)
    if ph.height > max_photo_h:
        ph.height = max_photo_h


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    image_candidates: tuple[str, ...],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _no_line(bg)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.name = "BGFill"

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.0))
    _no_line(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_BAR
    bar.name = "TitleBandSlide"

    ttl = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.2), SW - Inches(0.9), Inches(0.72)
    )
    ttl.name = "SlideHeading"
    tt = ttl.text_frame
    tt.clear()
    p0 = tt.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(25)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(255, 255, 255)

    body = slide.shapes.add_textbox(Inches(0.42), Inches(1.08), Inches(6.65), SH - Inches(1.3))
    body.name = "BulletStack"
    bf = body.text_frame
    bf.word_wrap = True
    bf.vertical_anchor = MSO_ANCHOR.TOP
    bf.auto_size = MSO_AUTO_SIZE.NONE
    bf.margin_left = Pt(4)
    bf.margin_right = Pt(6)
    bf.clear()
    shortened = bullets[:5]
    for i, raw in enumerate(shortened):
        line = raw.replace("`", "")
        if len(line) > 140:
            line = line[:137].rsplit(" ", 1)[0] + "…"
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = "• " + line
        para.font.size = Pt(13)
        para.font.color.rgb = C_TEXT
        para.space_after = Pt(6)

    img_path = resolve_img(*image_candidates)
    photo_top = Inches(1.05)
    photo_h = Inches(5.45)
    slide.shapes.add_picture(
        str(img_path),
        Inches(7.18),
        photo_top,
        height=photo_h,
    ).name = "SlidePhotoMain"

    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.05),
        Inches(0.92),
        Inches(6.16),
        Inches(5.72),
    )
    _no_line(frame)
    frame.fill.background()
    frame.line.color.rgb = C_FRAME
    frame.line.width = Pt(2)
    frame.name = "SlidePhotoAccent"


CONTENT: list[tuple[str, list[str], tuple[str, ...]]] = [
    (
        "Ansamblu — ce înglobează site-ul",
        [
            "Peste 50 simulări cu rută dedicată și filtre în catalog.",
            "Probleme + BAC + grile pe același cont; progres în cloud.",
            "Resurse /resurse pe capitole și avatar AI oriunde în SPA.",
            "Profil, realizări, comunitate; spațiu profesor–clasă.",
        ],
        ("Pendule_Multiple_Screenshot.png",),
    ),
    (
        "Pendule clasice și neliniare",
        [
            "Simplu, amortizat, neliniar, pendule Multiple → compară perioada.",
            "Grafice live legate clar de curricula de oscilații.",
        ],
        ("Simplu_Screenshot.png", "Amortizat_Screenshot.png"),
    ),
    (
        "Oscilații 1D și vizual grafică",
        [
            "OX / OY: amplitudine, fază pentru legături cu onde și bac.",
            "Construiești intuiția înaintea fig. Lissajous.",
        ],
        ("Oscilatieox_Screenshot.png",),
    ),
    (
        "Unde — propagare",
        [
            "Fronturi de unde, viteză λ-f; exemple din medii diverse.",
            "Leagă teoria capitole Unde în /resurse de laborator vizual.",
        ],
        ("Unde_Screenshot.png",),
    ),
    (
        "Figuri Lissajous",
        [
            "Raport frecvențe + defazaj → traiectorii clasice liceu.",
            "Ideal pentru itemi concurs și reprezentări bac.",
        ],
        ("Lissajous_Screenshot.png",),
    ),
    (
        "Seismologie",
        [
            "Propagare seismică și răspuns structural simplificat vizual.",
            "Context geofizică fără a părăsi browserul.",
        ],
        ("Seism_Screenshot.png",),
    ),
    (
        "Termodinamică — gaze și transformări",
        [
            "Simulator gaz ideal și diagrame stare.",
            "Baza pentru bac terminologie PV, lucru căldură.",
        ],
        ("Termodinamica_Screenshot.png",),
    ),
    (
        "Motoare termice",
        [
            "Otto, Diesel, Carnot comparate vizual pe același ecran.",
            "Formule bac + ciclu sincron între curs și sim.",
        ],
        ("Motoare_Termice_Screenshot.png",),
    ),
    (
        "Mecanică — proiectile, plan înclinat, coliziuni",
        [
            "Proiectile parabolice parametrizabile.",
            "Plan înclinat + coliziuni inelastice pentru mecanica liceu.",
        ],
        ("Proiectile_Screenshot.png", "Plan_Inclinat_Screenshot.png"),
    ),
    (
        "Termodinamică aplicată & transport",
        [
            "„Frecare aer’’ și exemple legate energeticii fluide fenomenologic.",
            "Susține explicația consumului rezistiv/atmosferă.",
        ],
        ("Frecare_Aer_Screenshot.png", "Motoare_Termice_Screenshot.png"),
    ),
    (
        "Optică — prismă, lentile, refracție",
        [
            "Prisme, lentilă subțire, refracție atmosferică pentru geometria raze.",
            "Reflexie–refracție în același modul grafic liceu.",
        ],
        ("Prisma_Screenshot.png", "Lentila_Subtire_Screenshot.png"),
    ),
    (
        "Lasere — propagare cavitate și fascicule",
        [
            "Simulator dedicat geometrie optică stimulated emission.",
            "Include variante grafice pentru emisie colimată liceu-modern.",
        ],
        ("Laser_Simulator_Screenshot.png", "laser_Screenshot.png"),
    ),
    (
        "Laseri mari — ELI‑NP și accelerator fotonic",
        [
            "Vizualizări infrastructură de tip ELI‑NP pentru context research.",
            "„Accelerator laser’’ — articularea dintre liceu și fotonica actuală.",
        ],
        ("Eli_Np_Laser_Screenshot.png", "Accelerator_Laser_Screenshot.png"),
    ),
    (
        "Electricitate & Kirchhoff",
        [
            "Montaje liceu: serie paralel energetici.",
            "Kirchhoff explicit pentru bac electricitate continuă.",
        ],
        ("Circuite_Electricitate_Screenshot.png", "kirchoff_Screenshot.png"),
    ),
    (
        "Curent alternativ & spectru",
        [
            "Oscilații sinusoidale parametrizabile și valorile RMS vizibile.",
            "Spectrul electromagnetic interactiv pentru lecție interdisciplinară.",
        ],
        ("ac_Screenshot.png", "spectru_Screenshot.png"),
    ),
    (
        "Interferență & exemple cuantice fenomenologice",
        [
            "Dublă fantă clasic experiment gând.",
            "Tunel cuantic pentru probabilistic intuitiv liceu-modern.",
        ],
        ("dubla_fanta_Screenshot.png", "tunelare_Screenshot.png"),
    ),
    (
        "Matematică pe grafici",
        [
            "Grafici funcție + pendul derivat pentru legături analiză liceu.",
            "Verificarea numerică paralelă rezolvărilor clasice.",
        ],
        ("Grafice_Basic_Screenshot.png", "Functii_Screenshot.png"),
    ),
    (
        "Michelson și fizica relativității (intro)",
        [
            "Interferometru istoric prezent grafic liceu-modern.",
            "Leagă undă lumină și experiența remarcabilă din manual.",
        ],
        ("Michaelson_Morley_Screenshot.png",),
    ),
    (
        "Astronomie & Kepler",
        [
            "Orbită și legile Kepler + constelații pentru capitol astro.",
            "Curriculum integrat liceu între cinematică și astro.",
        ],
        ("Miscare_Planete_Screenshot.png", "Legi_Kepler_Screenshot.png"),
    ),
    (
        "Atom, tabel periodic, legături",
        [
            "Hidrogen, tabloul periodic interactiv și legături atomice.",
            "Pod între curs chimie și fizica atomică.",
        ],
        ("Atom_Hidrogen_Screenshot.png", "Tabel_periodic_Screenshot.png"),
    ),
    (
        "Energie nucleară & izotopi",
        [
            "Fisiune vizual și reactoare educaționale; fuziune D–T.",
            "Familii de izotopi uraniu + catalog complet izotopi educativ.",
        ],
        (
            "Fisiune_Nucleara_Screenshot.png",
            "Reactor_Fuziune_Dt_Screenshot.png",
            "Izotopi_Uraniu_Screenshot.png",
        ),
    ),
    (
        "Apă grea & proces industrial lejer",
        [
            "Model schimb izotopic și apă grea.",
            "Distilație fracționată prezentată ca resursă vizual tehnic.",
        ],
        (
            "apa_grea_1.png",
            "schimb_izotopic_Screenshot.png",
            "Distilare_D2o_Fractionata_Resurse.png",
        ),
    ),
    (
        "Hidrogen criogenic & supraconductibilitate",
        [
            "Criogenie, supraconductivitate, celule combustibil illustrative.",
            "Extinde platforma spre tehnică de laborator avansată.",
        ],
        (
            "Criogenie_Screenshot.png",
            "Supraconductivitate_Screenshot.png",
            "Fuel_Cell_Screenshot.png",
        ),
    ),
    (
        "Probleme, BAC și grile",
        [
            "Catalog filtrabil; pagini /probleme/:id rezolvabile.",
            "Secțiuni /probleme/bac și /probleme/grile pentru repetare eficientă.",
            "Feedback automat + API Puls AI pentru fizică pe înțeles.",
        ],
        ("Prisma_Screenshot.png", "Grafice_Basic_Screenshot.png"),
    ),
    (
        "Resurse capitole multimedia",
        [
            "/resurse + subrute capitole Pendule→Lasere→Nuclear.",
            "Embed video/MathJax oriunde apare teoria textuală.",
        ],
        ("Constelatii_Screenshot.png", "Michaelson_Morley_Screenshot.png"),
    ),
    (
        "Asistent AI și Puls AI",
        [
            "Avatar Three.js oriunde navighezi; flux n8n + context site.",
            "Puls AI separat pentru rezolvări BAC și audit teme fizică.",
        ],
        ("Vizualizator_4d_Screenshot.png", "Laser_Simulator_Screenshot.png"),
    ),
    (
        "Gamificare, profil și comunitate",
        [
            "Badge-uri, statistici, profil /profil/:alias dacă vizibil.",
            "/comunitate pentru socializarea progresului educațional.",
        ],
        ("pendule_placeholder",),
    ),
    (
        "Profesor, clase și admin",
        [
            "/profesor dashboard + clase; elevi intra cu cod/link.",
            "/admin pentru întreținere conținut (Firebase roles).",
        ],
        ("pendule_placeholder",),
    ),
    (
        "Limbi și experiență UI",
        [
            "RO + EN pentru public internațional și diaspora.",
            "Dark/light, căutare, layout responsive pentru telefon liceu.",
        ],
        ("pendule_placeholder",),
    ),
    (
        "Arhitectură tehnică (rezumat)",
        [
            "React 19, Vite, SCSS modular, Redux toolkit.",
            "Firebase datastore; CDN imagini Cloudinary/ImageKit;",
            "Vercel deploy; automatizare n8n + Groq Llama pentru chat.",
        ],
        ("kirchoff_Screenshot.png", "ac_Screenshot.png"),
    ),
    (
        "Concluzie",
        [
            "PULS livrează tot parcursul: vizual → teorie → practică bac.",
            "Echipă liceu ghidată; mentori fizică pentru corectitudine științifică.",
            "Descoperă, experimentează, pulsul fascinant al științei!",
        ],
        ("Lissajous_Screenshot.png", "Unde_Screenshot.png"),
    ),
]


def _materialize_fallbacks(
    rows: list[tuple[str, list[str], tuple[str, ...]]],
    fallback_pngs: tuple[str, ...],
) -> list[tuple[str, list[str], tuple[str, ...]]]:
    fb = fallback_pngs
    out = []
    for title, bulls, imgs in rows:
        if imgs and imgs[0] == "pendule_placeholder":
            imgs = fb
        out.append((title, bulls, imgs))
    return out


CONTENT = _materialize_fallbacks(
    CONTENT,
    (
        "Grafice_Basic_Screenshot.png",
        "Functii_Screenshot.png",
    ),
)


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    add_title_slide(prs)
    for title, bulls, imgs in CONTENT:
        add_content_slide(prs, title, bulls, imgs)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Întrebări?"

    target = OUT
    try:
        save_presentation_with_morph(prs, target)
        print(f"Saved (Morph injectat în OOXML): {target}")
    except PermissionError:
        alt = OUT.with_name(f"{OUT.stem}-nou{OUT.suffix}")
        save_presentation_with_morph(prs, alt)
        print(f"{OUT.name} blocat; Saved: {alt}")


if __name__ == "__main__":
    main()
